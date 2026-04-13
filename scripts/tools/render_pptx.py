"""
render_pptx.py
JSON または Markdown から .pptx を生成する簡易実装（python-pptx が必要）
"""

import argparse
import json
import re
import sys
from pathlib import Path


def parse_markdown_to_slides(md: str):
    parts = [s.strip() for s in md.split('\n---\n') if s.strip()]
    slides = []
    for i, sec in enumerate(parts):
        lines = sec.splitlines()
        title = lines[0].lstrip('# ').strip() if lines and lines[0].startswith('#') else f"Slide {i+1}"
        content = [ln.strip('- ').strip() for ln in lines[1:] if ln.strip()]
        slides.append({"layout": "content", "title": title, "content": content})
    return {"title": "Presentation", "slides": slides}


def render_slides_to_pptx(slides_data: dict, output_path: str) -> dict:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except Exception as e:
        return {"success": False, "error": f"python-pptx not available: {e}", "output_path": None}

    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    # title slide
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1] if len(slide.placeholders) > 1 else None
    title.text = slides_data.get("title", "Presentation")
    if subtitle and slides_data.get("subtitle"):
        subtitle.text = slides_data.get("subtitle")

    for s in slides_data.get("slides", []):
        layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)
        if slide.shapes.title:
            slide.shapes.title.text = s.get("title", "")
        body_shape = None
        for shp in slide.shapes:
            if shp.has_text_frame:
                body_shape = shp
                break
        if body_shape is None:
            continue
        tf = body_shape.text_frame
        tf.clear()
        for i, line in enumerate(s.get("content", [])):
            if i == 0:
                tf.text = line
            else:
                p = tf.add_paragraph()
                p.text = line
                p.level = 0

    outp = Path(output_path)
    prs.save(str(outp))
    size = outp.stat().st_size if outp.exists() else 0
    return {"success": True, "output_path": str(outp), "size_bytes": size}


def load_slides_data(inp: Path) -> dict:
    """ファイルを読み込んでスライドデータを返す。JSON・Markdown・JSONブロック混在に対応"""
    text = inp.read_text(encoding="utf-8").strip()

    # ```json ... ``` ブロックを抽出
    json_block = re.search(r'```(?:json)?\s*(\{.*?\})\s*```', text, re.DOTALL)
    if json_block:
        text = json_block.group(1).strip()

    # JSON として解析を試みる
    if inp.suffix.lower() == ".json" or text.startswith("{"):
        try:
            return json.loads(text)
        except json.JSONDecodeError:
            pass  # JSON失敗 → Markdownとして解析

    # Markdown として解析
    return parse_markdown_to_slides(text)


if __name__ == "__main__":
    parser = argparse.ArgumentParser()
    parser.add_argument("--input", required=True)
    parser.add_argument("--output", required=True)
    args = parser.parse_args()

    inp = Path(args.input)
    if not inp.exists():
        print(json.dumps({"success": False, "error": "input not found"}, ensure_ascii=False))
        sys.exit(1)

    try:
        data = load_slides_data(inp)
    except Exception as e:
        print(json.dumps({"success": False, "error": f"parse error: {e}"}, ensure_ascii=False))
        sys.exit(1)

    res = render_slides_to_pptx(data, args.output)
    print(json.dumps(res, ensure_ascii=False))
    sys.exit(0 if res.get("success") else 1)
